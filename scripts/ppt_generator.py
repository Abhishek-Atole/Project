from pptx import Presentation

def generate_ppt(mcq_data, output_path):
    """
    Generate a PowerPoint presentation from MCQ data.
    """
    print(f"Generating PowerPoint at: {output_path}")  # Debugging statement
    prs = Presentation()

    for question_data in mcq_data["questions"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

        # Set the question as the slide title
        slide.shapes.title.text = question_data["question"]

        # Format the choices as bullet points
        choices_text = "\n".join(question_data["choices"])
        slide.placeholders[1].text = choices_text

        # Add images if available
        for image_path in question_data["images"]:
            slide.shapes.add_picture(image_path, left=100, top=200, width=300, height=200)

    # Save the PowerPoint file
    prs.save(output_path)
    print(f"PowerPoint saved successfully at: {output_path}")  # Debugging statement